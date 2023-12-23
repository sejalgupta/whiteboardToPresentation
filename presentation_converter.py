from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from constants import layout_mapping
import json
from helper_functions import scale_image
from presentation_theme import PresentationTheme
from slide_content import get_metadata_gpt
from theme import get_theme_gpt

def add_slide(presentation, information):
    """
    Add a slide to the presentation with flexible content handling.

    Parameters
    - presentation (Presentation obj): the main presentation that is being modified
    - information (dictionary): the layout of the slide and data being placed on the slide including the titles, subtitles, content, and more
    """

    layout = layout_mapping[information["layout"]]
    slide_layout = presentation.slide_layouts[layout.value]
    slide = presentation.slides.add_slide(slide_layout)

    for placeholder in slide.placeholders:
        print(placeholder.name)
        # ADD SUBTITLE
        if ('subtitle' in placeholder.name.lower() or "text" in placeholder.name.lower()) and len(information["subtitle"]) > 0:
            title = information["subtitle"].pop(0)
            placeholder.text = title
        
        # ADD TITLE
        elif 'title' in placeholder.name.lower() and len(information["title"]) > 0:
            title = information["title"].pop(0)
            placeholder.text = title

        # ADD PICTURE / ICON
        elif 'picture' in placeholder.name.lower() and len(information["picture"]) > 0:
            img_path = information["picture"].pop(0)
            placeholder.insert_picture(img_path)

        # ADD IMAGE OR TEXT IN CONTENT HOLDER
        elif 'content' in placeholder.name.lower() and len(information["content"]) > 0:
            content_dict = information["content"].pop(0)
            if content_dict["type"] == "text":
                tf = placeholder.text_frame
                tf.text = content_dict["text"]
            else:
                img_path = content_dict["image"]
                image_width, image_height = scale_image(img_path, placeholder.width, placeholder.height)

                # add image as a shape
                slide.shapes.add_picture(img_path, placeholder.left, placeholder.top, image_width, image_height)

def create_presentation(name, role, presentation_context, presentation_length):
    pres = Presentation()

    # #load theme content
    # theme_content = None
    # try:
    #     theme_content = json.load(open(f"data.json"))
    # except:
    #     answers = get_theme_gpt(role, presentation_context)
    #     response = answers["choices"][0]["message"]["content"]

    #     theme_content = json.loads(response)
    #     with open('data.json', 'w') as f:
    #         json.dump(response, f)
        
    # theme = PresentationTheme(theme_content)

    # print(theme.get_theme_info())

    #load slides content
    slides_content = None
    try:
        slides_content = json.load(open(f"json/{name}.json"))
    except:
        answers = get_metadata_gpt(role, presentation_context, presentation_length)
        response = answers["choices"][0]["message"]["content"]

        slides_content = json.loads(response)
        with open(f"json/{name}.json", 'w') as f:
            json.dump(slides_content, f)

    #add each slide
    for slide in slides_content["slides"]:
        try:
            add_slide(pres, slide)
        except:
            print("ERROR", slide["layout"])

    pres.save(f'presentation/{name}.pptx')
    print("Presentation created successfully:", name)

def main():
    for name, role, presentation_context, presentation_length in [
        ("multiplying_fractions", "teacher for sixth grade math", "I am trying to teach about multiplying fractions. I need to create a lesson plan for tomorrow.", "10-minute"),
        ("dividing_fractions", "teacher for sixth grade math", "I am trying to teach about dividing fractions. I need to create a lesson plan for tomorrow.", "10-minute")
    ]:
        create_presentation(name, role, presentation_context, presentation_length)

if __name__ == "__main__":
    main()