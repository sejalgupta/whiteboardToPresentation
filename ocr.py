
from google.cloud import vision
from pptx import Presentation 
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

"""
@params: path- image_path
         slide- the slide you want to add the image on 

For each box that it recognizes, creates a bounding box using the scaled metrics
and adds that as a rectangular shape with the corresponding text in a slide 
 
"""
def detect_document(path,slide):
    """Detects document features in an image."""
    client = vision.ImageAnnotatorClient()

    # get image dimensions 
    with Image.open(path) as img:
        width, height = img.size
    

    with open(path, "rb") as image_file:
        content = image_file.read()

    image = vision.Image(content=content)

    response = client.document_text_detection(image=image)

    for page in response.full_text_annotation.pages:
        for block in page.blocks:
            #print(f"\nBlock confidence: {block.confidence}\n")

            for paragraph in block.paragraphs:
                #print("Paragraph confidence: {}".format(paragraph.confidence))

                for word in paragraph.words:
                    
                    word_text = "".join([symbol.text for symbol in word.symbols])
                    #print(word_text)
                    X_SCALE = 10/width
                    Y_SCALE =  7.5/height
                    # get bounding box of the symbols 
                    left_top = word.symbols[0].bounding_box.vertices[0]
                    left_bottom = word.symbols[0].bounding_box.vertices[3]
                    right_top = word.symbols[-1].bounding_box.vertices[1]
                    right_bottom = word.symbols[-1].bounding_box.vertices[2]
                    #bounding box 
                    w = (right_top.x)*X_SCALE-(left_top.x)*X_SCALE
                    h = (right_bottom.y)*Y_SCALE-(right_top.y)*Y_SCALE 
                    # bounding box 
                    # scale each one 
                    # Format: [left, top, width, height] in inches
                    bounding_box = [left_top.x*X_SCALE, left_top.y*Y_SCALE, w, h]  # Example coordinates
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,  # Using rectangle shape for the square
                        Inches(bounding_box[0]),
                        Inches(bounding_box[1]),
                        Inches(bounding_box[2]),
                        Inches(bounding_box[3])
                    )
                    shape.text = word_text
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)  # Set font size to 12 point
                            run.font.name = 'Arial'  # Set font type to Arial
                            run.font.bold = True  # Make text bold
                            run.font.italic = False  # Italicize text (set to True if needed)
                            run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black

                    
                    """
                    print(
                        "Word text: {} (confidence: {})".format(
                            word_text, word.confidence
                        )
                    )
                    """
                    """
                    for symbol in word.symbols:
                        vertices = [(v.x, v.y) for v in symbol.bounding_box.vertices]
                        print(
                            "\tSymbol: {} (confidence: {}) - Coordinates: {}".format(
                                symbol.text, symbol.confidence, vertices
                            )
                        )
                    """

    if response.error.message:
        raise Exception(
            "{}\nFor more info on error messages, check: "
            "https://cloud.google.com/apis/design/errors".format(response.error.message)
        )


    
if __name__ == "__main__":
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Using blank layout
    slide = presentation.slides.add_slide(slide_layout)
    detect_document('handwritten.png',slide)
    presentation_file = 'presentation_with_square.pptx'
    presentation.save(presentation_file) 
